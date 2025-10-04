#!/usr/bin/env python3
"""Extract text from a .pptx PowerPoint file"""
import os
from pptx import Presentation


def extract_table_text(table):
    """Yield each row of a table as a formatted string."""
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
        if cells:
            yield " | ".join(cells)


def extract_text(pptx_path):
    """
    Yield text chunks slide-by-slide from a PowerPoint presentation.

    Args:
        pptx_path (str): Path to the .pptx file
    """
    if not pptx_path.lower().endswith(".pptx"):
        raise ValueError("Input file must be a .pptx file")

    prs = Presentation(pptx_path)

    for slide_num, slide in enumerate(prs.slides, start=1):
        content_parts = [f"--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if not shape.has_text_frame and not shape.has_table:
                continue

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        content_parts.append(text)

            if shape.has_table:
                for line in extract_table_text(shape.table):
                    content_parts.append(line)

        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                note_text = notes_slide.notes_text_frame.text.strip()
                if note_text:
                    content_parts.append(f"[Notes] {note_text}")

        slide_text = "\n".join(content_parts).strip()
        if slide_text:
            yield f"{slide_text}\n(Slide {slide_num})"
